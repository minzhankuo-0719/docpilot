"""PPTX parsing — each paragraph becomes a Block.

Title placeholders are tagged ``block_type="heading"`` so the chunker can
treat slide titles as standalone retrieval units.
"""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation

from .pdf import Block


@dataclass
class ParsedSlide:
    slide_num: int                      # 1-indexed
    text: str                           # blocks joined by "\n\n" (back-compat)
    source: str
    blocks: list[Block] = field(default_factory=list)


def _is_title_shape(shape) -> bool:
    if not shape.is_placeholder:
        return False
    fmt = shape.placeholder_format
    if fmt is None:
        return False
    # idx == 0 is the slide title placeholder per OOXML spec.
    return fmt.idx == 0


def parse_pptx(path: str | Path) -> list[ParsedSlide]:
    """Extract text from every slide of a PPTX, preserving paragraph blocks."""
    path = Path(path)
    source = path.name
    prs = Presentation(str(path))
    slides: list[ParsedSlide] = []

    for i, slide in enumerate(prs.slides, 1):
        blocks: list[Block] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            is_title = _is_title_shape(shape)
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if not line:
                    continue
                blocks.append(Block(
                    text=line,
                    block_type="heading" if is_title else "paragraph",
                ))

        joined = "\n\n".join(b.text for b in blocks)
        slides.append(ParsedSlide(
            slide_num=i,
            text=joined,
            source=source,
            blocks=blocks,
        ))

    return slides
